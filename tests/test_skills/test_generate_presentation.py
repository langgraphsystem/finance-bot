"""Tests for generate_presentation skill."""

from unittest.mock import AsyncMock, patch

from src.gateway.types import IncomingMessage, MessageType
from src.skills.generate_presentation.handler import skill


async def test_generate_presentation_no_topic(sample_context):
    """No topic provided — asks user what to create."""
    msg = IncomingMessage(
        id="msg-1",
        user_id="tg_1",
        chat_id="chat_1",
        type=MessageType.text,
        text="",
    )
    result = await skill.execute(msg, sample_context, {})
    # Response may be in Russian (sample_context.language="ru") or English
    assert result.response_text  # Non-empty prompt asking about the topic
    assert result.document is None


async def test_generate_presentation_attributes():
    assert skill.name == "generate_presentation"
    assert "generate_presentation" in skill.intents
    assert skill.model == "claude-sonnet-4-6"


async def test_generate_presentation_e2b_happy_path(sample_context):
    """E2B sandbox succeeds — returns presentation document."""
    msg = IncomingMessage(
        id="msg-1",
        user_id="tg_1",
        chat_id="chat_1",
        type=MessageType.text,
        text="create a presentation about AI trends",
    )
    pptx_bytes = b"PK\x03\x04fake-pptx-output"
    with (
        patch(
            "src.skills.generate_presentation.handler.generate_text",
            new_callable=AsyncMock,
            return_value=(
                "from pptx import Presentation\nprs=Presentation()\nprs.save('/tmp/output.pptx')"
            ),
        ),
        patch(
            "src.tools.e2b_file_utils.execute_code_with_file",
            new_callable=AsyncMock,
            return_value=(pptx_bytes, "Success"),
        ),
    ):
        result = await skill.execute(
            msg, sample_context, {"presentation_topic": "AI trends in 2026"}
        )

    assert result.document == pptx_bytes
    assert result.document_name.endswith(".pptx")
    assert "ready" in result.response_text.lower()


async def test_generate_presentation_e2b_fails_fallback_succeeds(sample_context):
    """E2B fails — falls back to local python-pptx JSON spec generation."""
    msg = IncomingMessage(
        id="msg-1",
        user_id="tg_1",
        chat_id="chat_1",
        type=MessageType.text,
        text="create a sales pitch presentation",
    )
    fallback_pptx = b"PK\x03\x04local-pptx-output"
    with (
        patch(
            "src.skills.generate_presentation.handler.generate_text",
            new_callable=AsyncMock,
            return_value=(
                "from pptx import Presentation\nprs=Presentation()\nprs.save('/tmp/output.pptx')"
            ),
        ),
        patch(
            "src.tools.e2b_file_utils.execute_code_with_file",
            new_callable=AsyncMock,
            side_effect=Exception("E2B unavailable"),
        ),
        patch(
            "src.skills.generate_presentation.handler._build_fallback_pptx",
            new_callable=AsyncMock,
            return_value=fallback_pptx,
        ),
    ):
        result = await skill.execute(msg, sample_context, {"presentation_topic": "sales pitch"})

    assert result.document == fallback_pptx
    assert result.document_name.endswith(".pptx")
    assert "ready" in result.response_text.lower()


async def test_generate_presentation_all_methods_fail(sample_context):
    """Both E2B and fallback fail — returns error message."""
    msg = IncomingMessage(
        id="msg-1",
        user_id="tg_1",
        chat_id="chat_1",
        type=MessageType.text,
        text="create a presentation about quantum computing",
    )
    with (
        patch(
            "src.skills.generate_presentation.handler.generate_text",
            new_callable=AsyncMock,
            return_value="code here",
        ),
        patch(
            "src.tools.e2b_file_utils.execute_code_with_file",
            new_callable=AsyncMock,
            side_effect=Exception("E2B down"),
        ),
        patch(
            "src.skills.generate_presentation.handler._build_fallback_pptx",
            new_callable=AsyncMock,
            return_value=None,
        ),
    ):
        result = await skill.execute(
            msg, sample_context, {"presentation_topic": "quantum computing"}
        )

    assert "failed" in result.response_text.lower() or "different" in result.response_text.lower()
    assert result.document is None


async def test_generate_presentation_uses_message_text(sample_context):
    """No presentation_topic in intent_data — falls back to message text."""
    msg = IncomingMessage(
        id="msg-1",
        user_id="tg_1",
        chat_id="chat_1",
        type=MessageType.text,
        text="create presentation about marketing strategy",
    )
    pptx_bytes = b"PK\x03\x04pptx-output"
    with (
        patch(
            "src.skills.generate_presentation.handler.generate_text",
            new_callable=AsyncMock,
            return_value=(
                "from pptx import Presentation\nprs=Presentation()\nprs.save('/tmp/output.pptx')"
            ),
        ),
        patch(
            "src.tools.e2b_file_utils.execute_code_with_file",
            new_callable=AsyncMock,
            return_value=(pptx_bytes, "OK"),
        ),
    ):
        result = await skill.execute(msg, sample_context, {})

    assert result.document == pptx_bytes
    assert "marketing" in result.document_name.lower()
