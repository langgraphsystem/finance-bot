"""Generate PPTX presentations via LLM-authored python-pptx code in E2B sandbox."""

import asyncio
import io
import logging
import re
from typing import Any

from src.core.context import SessionContext
from src.core.llm.clients import generate_text
from src.core.observability import observe
from src.gateway.types import IncomingMessage
from src.skills.base import SkillResult

logger = logging.getLogger(__name__)

PPTX_SYSTEM_PROMPT = """\
You are a presentation generator. You write Python code using python-pptx
that creates a .pptx file at /tmp/output.pptx.

Rules:
- Always start with: from pptx import Presentation
- Always save to /tmp/output.pptx
- Use pptx.util (Inches, Pt, Emu) for sizing
- Use pptx.dml.color.RGBColor for colors
- Use pptx.enum.text.PP_ALIGN for alignment
- Create a professional presentation with:
  * Title slide with topic and subtitle
  * Content slides with bullet points (3-6 slides total)
  * Clean layout, readable fonts (18-44pt)
  * Consistent color scheme
- Add real, substantive content — not placeholder text
- The code must run standalone without errors
- Do NOT use any library except python-pptx and its submodules
- Respond ONLY with Python code, no explanations outside the code"""

FALLBACK_SYSTEM_PROMPT = """\
You are a presentation generator. Describe the presentation structure as JSON.
Return a JSON object with:
{
  "title": "Presentation title",
  "subtitle": "Optional subtitle",
  "slides": [
    {
      "title": "Slide title",
      "bullets": ["Point 1", "Point 2", "Point 3"]
    }
  ]
}
Create 4-8 slides with substantive content. Respond ONLY with valid JSON."""


def _strip_markdown_fences(text: str) -> str:
    """Remove ```language ... ``` wrappers from LLM output."""
    text = text.strip()
    match = re.match(r"^```\w*\n(.*?)```$", text, re.DOTALL)
    if match:
        return match.group(1).strip()
    return text


async def _build_fallback_pptx(description: str) -> bytes | None:
    """Generate a presentation without E2B using LLM JSON + local python-pptx."""
    import json

    spec_text = await generate_text(
        model="claude-sonnet-4-6",
        system=FALLBACK_SYSTEM_PROMPT,
        prompt=f"Create a presentation about: {description}",
        max_tokens=2048,
    )
    spec_text = _strip_markdown_fences(spec_text)

    try:
        spec = json.loads(spec_text)
    except json.JSONDecodeError:
        logger.warning("Failed to parse presentation JSON spec")
        return None

    def _create():
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = spec.get("title", "Presentation")
        if slide.placeholders[1]:
            slide.placeholders[1].text = spec.get("subtitle", "")

        # Content slides
        bullet_layout = prs.slide_layouts[1]
        for slide_data in spec.get("slides", []):
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = slide_data.get("title", "")

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            bullets = slide_data.get("bullets", [])
            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.level = 0

            # Style bullets
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                paragraph.alignment = PP_ALIGN.LEFT

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    return await asyncio.to_thread(_create)


class GeneratePresentationSkill:
    name = "generate_presentation"
    intents = ["generate_presentation"]
    model = "claude-sonnet-4-6"

    @observe(name="generate_presentation")
    async def execute(
        self,
        message: IncomingMessage,
        context: SessionContext,
        intent_data: dict[str, Any],
    ) -> SkillResult:
        topic = (intent_data.get("presentation_topic") or message.text or "").strip()

        if not topic:
            lang = context.language or "en"
            if lang == "ru":
                prompt = "О чём должна быть презентация?"
            elif lang == "es":
                prompt = "Sobre que debe ser la presentacion?"
            else:
                prompt = "What should the presentation be about?"
            return SkillResult(response_text=prompt)

        # Try E2B sandbox first
        try:
            from src.tools.e2b_file_utils import execute_code_with_file

            code = await generate_text(
                model=self.model,
                system=PPTX_SYSTEM_PROMPT,
                prompt=f"Create a presentation about: {topic}",
                max_tokens=4096,
            )
            code = _strip_markdown_fences(code)

            file_bytes, stdout = await execute_code_with_file(
                code=code,
                output_filename="output.pptx",
                language="python",
                timeout=60,
                install_deps=["python-pptx"],
            )

            if file_bytes:
                filename = _make_filename(topic)
                logger.info(
                    "Presentation generated via E2B for user %s (%d bytes)",
                    context.user_id,
                    len(file_bytes),
                )
                return SkillResult(
                    response_text=f"<b>{filename}</b> — your presentation is ready.",
                    document=file_bytes,
                    document_name=filename,
                )

            logger.warning("E2B produced no file, falling back to local: %s", stdout)
        except Exception as e:
            logger.warning("E2B presentation generation failed, using fallback: %s", e)

        # Fallback: local python-pptx via LLM JSON spec
        file_bytes = await _build_fallback_pptx(topic)
        if file_bytes:
            filename = _make_filename(topic)
            logger.info(
                "Presentation generated locally for user %s (%d bytes)",
                context.user_id,
                len(file_bytes),
            )
            return SkillResult(
                response_text=f"<b>{filename}</b> — your presentation is ready.",
                document=file_bytes,
                document_name=filename,
            )

        return SkillResult(response_text="Failed to generate presentation. Try a different topic.")

    def get_system_prompt(self, context: SessionContext) -> str:
        return PPTX_SYSTEM_PROMPT


def _make_filename(topic: str) -> str:
    """Generate a short filename from the topic."""
    slug = re.sub(r"[^a-zA-Z0-9\s]", "", topic.lower())
    slug = "_".join(slug.split()[:5])
    if not slug:
        slug = "presentation"
    if len(slug) > 40:
        slug = slug[:40].rstrip("_")
    return f"{slug}.pptx"


skill = GeneratePresentationSkill()
